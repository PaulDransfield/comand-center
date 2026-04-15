"""
Presentation Generator — Restaurang Björken
Generates real .pptx files using python-pptx.
Aesthetic: Midnight Executive — deep navy dominates, ice-blue accents, white type.
"""
import sys, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── COLOUR THEMES ────────────────────────────────────────────────
THEMES = {
    "midnight": {
        "name":    "Midnight Executive",
        "bg_dark": RGBColor(0x1E, 0x27, 0x61),   # deep navy
        "bg_mid":  RGBColor(0x13, 0x1A, 0x45),   # darker navy
        "bg_light":RGBColor(0xCA, 0xDC, 0xFC),   # ice blue
        "accent":  RGBColor(0xCA, 0xDC, 0xFC),   # ice blue
        "accent2": RGBColor(0xFF, 0xC0, 0x40),   # gold
        "white":   RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark":RGBColor(0x1E,0x27,0x61),
        "text_muted":RGBColor(0x8A,0x9B,0xC8),
        "chart_colors": ["CADCFC","FFC040","4A7FC1","7EC8E3","E8EDF8"],
    },
    "forest": {
        "name":    "Forest & Moss",
        "bg_dark": RGBColor(0x2C, 0x5F, 0x2D),
        "bg_mid":  RGBColor(0x1A, 0x3A, 0x1B),
        "bg_light":RGBColor(0x97, 0xBC, 0x62),
        "accent":  RGBColor(0x97, 0xBC, 0x62),
        "accent2": RGBColor(0xF5, 0xC5, 0x18),
        "white":   RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark":RGBColor(0x2C,0x5F,0x2D),
        "text_muted":RGBColor(0x7A,0xA8,0x5A),
        "chart_colors": ["97BC62","2C5F2D","F5C518","4A7A45","C8E6A0"],
    },
    "terracotta": {
        "name":    "Warm Terracotta",
        "bg_dark": RGBColor(0xB8, 0x50, 0x42),
        "bg_mid":  RGBColor(0x7A, 0x32, 0x28),
        "bg_light":RGBColor(0xE7, 0xE8, 0xD1),
        "accent":  RGBColor(0xA7, 0xBE, 0xAE),
        "accent2": RGBColor(0xF5, 0xC5, 0x18),
        "white":   RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark":RGBColor(0x7A,0x32,0x28),
        "text_muted":RGBColor(0xB8,0x8A,0x80),
        "chart_colors": ["B85042","A7BEAE","E7E8D1","D4856B","7A3228"],
    },
}

# ── FINANCIAL DATA ────────────────────────────────────────────────
def get_data(overrides=None):
    d = {
        "company":      "Restaurang Björken AB",
        "period":       "Mars 2026",
        "quarter":      "Q1 2026",
        "revenue":      505900,
        "prev_revenue": 461800,
        "staff":        224978, "staff_pct": 44.5,
        "food":         149300, "food_pct":  29.5,
        "rent":          64500, "rent_pct":  12.7,
        "other":         28422, "other_pct":  5.6,
        "budget_revenue": 490000,
        "monthly_revenue": [389200, 461800, 505900],
        "monthly_profit":  [12100,  38700,  38700],
        "monthly_labels":  ["Jan", "Feb", "Mars"],
        "q1_revenue": 1356900, "q1_budget": 1350000,
    }
    d["total_costs"] = d["staff"] + d["food"] + d["rent"] + d["other"]
    d["profit"]      = d["revenue"] - d["total_costs"]
    d["margin"]      = round(d["profit"] / d["revenue"] * 100, 2)
    if overrides:
        d.update(overrides)
    return d

def kr(n):  return f"{int(abs(n)):,}".replace(",", " ") + " kr"
def pct(n): return f"{float(n):.1f}%"

# ── SLIDE BUILDER HELPERS ─────────────────────────────────────────
W, H = Inches(10), Inches(5.625)   # 16:9

def rgb(r,g,b): return RGBColor(r,g,b)

def add_rect(slide, x, y, w, h, fill_color, transparency=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.line.color.rgb = fill_color
    shape.line.width = Pt(0)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if transparency:
        shape.fill.fore_color.alpha = int(255 * (1 - transparency/100))
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color: run.font.color.rgb = color
    return tb

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_divider(slide, y, color, width=0.04, length=9.0, x_start=0.5):
    add_rect(slide, x_start, y, length, width, color)

def add_kpi_card(slide, x, y, w, h, value, label, t, color=None):
    """A KPI card: coloured rect with large value + small label"""
    c = color or t["accent"]
    # card background — use white with left accent bar
    add_rect(slide, x, y, w, h, rgb(0xFF,0xFF,0xFF))
    add_rect(slide, x, y, 0.06, h, c)  # accent bar
    add_text(slide, value, x+0.15, y+0.06, w-0.2, h*0.55,
             size=22, bold=True, color=t["text_dark"], font="Calibri")
    add_text(slide, label, x+0.15, y+h*0.58, w-0.2, h*0.38,
             size=10, color=t["text_muted"], font="Calibri")

def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text

# ── SLIDE CONSTRUCTORS ────────────────────────────────────────────

def slide_title(prs, t, d, title, subtitle, deck_type="Board Meeting"):
    """Slide 1: Full-bleed dark title slide"""
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(s, t["bg_dark"])

    # Accent band top
    add_rect(s, 0, 0, 10, 0.08, t["accent"])
    # Accent band bottom
    add_rect(s, 0, 5.545, 10, 0.08, t["accent"])

    # Left accent bar
    add_rect(s, 0.5, 1.2, 0.08, 3.2, t["accent2"])

    # Deck type label
    add_text(s, deck_type.upper(), 0.75, 1.25, 8, 0.45,
             size=11, color=t["accent2"], bold=True, font="Calibri",
             align=PP_ALIGN.LEFT)

    # Main title
    add_text(s, title, 0.75, 1.7, 8.5, 1.6,
             size=40, bold=True, color=t["white"], font="Calibri")

    # Subtitle
    add_text(s, subtitle, 0.75, 3.35, 8, 0.6,
             size=18, color=t["accent"], italic=True, font="Calibri")

    # Company + date bottom-left
    add_text(s, f"{d['company']}  ·  {d['period']}", 0.5, 5.1, 6, 0.35,
             size=11, color=t["text_muted"], font="Calibri")

    add_notes(s, f"Opening slide for {deck_type}.\n\nWelcome attendees. State the purpose: review {d['period']} performance and agree on Q2 actions. Expected duration: 45 minutes.")
    return s

def slide_exec_summary(prs, t, d):
    """Slide 2: Executive summary with 4 KPI cards"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    # Header bar
    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Executive Summary", 0.5, 0.17, 7, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")
    add_text(s, d["period"], 8.2, 0.22, 1.6, 0.46,
             size=13, color=t["accent"], font="Calibri", align=PP_ALIGN.RIGHT)

    # 4 KPI cards
    cards = [
        (kr(d["revenue"]),      "Revenue",      t["accent"]),
        (kr(d["total_costs"]),  "Total Costs",  rgb(0xE2,0x4B,0x4A)),
        (kr(d["profit"]),       "Net Profit",   rgb(0x52,0xC3,0x74)),
        (f"{d['margin']}%",     "Margin",       rgb(0x52,0xC3,0x74)),
    ]
    for i, (val, lbl, col) in enumerate(cards):
        add_kpi_card(s, 0.35 + i*2.35, 1.1, 2.2, 1.1, val, lbl, t, col)

    # Summary bullets
    bullets = [
        f"Revenue {kr(d['revenue'])} — up {round((d['revenue']-d['prev_revenue'])/d['prev_revenue']*100,1)}% vs February",
        f"Staff costs at {pct(d['staff_pct'])} of revenue — above 40% target; high sick leave in weeks 11–12",
        f"Food & Bev costs {pct(d['food_pct'])} — Menigo price increase +8% from March 1st",
        f"Net profit {kr(d['profit'])} at {pct(d['margin'])} margin — pressured by cost increases",
        "Q1 2026 revenue forecast exceeds budget by +6 900 kr",
    ]
    for i, b in enumerate(bullets):
        add_rect(s, 0.35, 2.4 + i*0.58, 0.06, 0.34, t["bg_dark"])
        add_text(s, b, 0.55, 2.38 + i*0.58, 9.1, 0.4,
                 size=13, color=rgb(0x2A,0x2A,0x2A), font="Calibri")

    add_notes(s, "Key talking points:\n• Revenue growth is strong but costs are outpacing it\n• Staff cost overrun is the single biggest issue — driven by sick leave and overtime\n• Menigo price increase adds approximately 15 000 kr/month to food costs\n• Action needed: renegotiate Menigo, review staffing schedule")
    return s

def slide_revenue(prs, t, d):
    """Slide 3: Revenue analysis with bar chart"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Revenue Analysis", 0.5, 0.17, 7, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")

    # Bar chart — monthly revenue
    chart_data = ChartData()
    chart_data.categories = d["monthly_labels"]
    chart_data.add_series("Revenue (kr)", d["monthly_revenue"])
    chart_data.add_series("Budget (kr)",  [380000, 460000, 490000])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(1.1), Inches(5.8), Inches(3.9),
        chart_data
    ).chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = t["bg_dark"]
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = t["accent"]

    # Right panel — breakdown
    breakdown = [
        ("Mat",      "61.8%", kr(312500)),
        ("Dryck",    "34.6%", kr(175000)),
        ("Catering", "3.6%",  kr(18400)),
    ]
    add_text(s, "Revenue breakdown", 6.5, 1.1, 3.2, 0.45,
             size=14, bold=True, color=t["text_dark"], font="Calibri")
    for i, (cat, p, amt) in enumerate(breakdown):
        y = 1.65 + i * 1.18
        add_rect(s, 6.5, y, 3.2, 1.0, rgb(0xFF,0xFF,0xFF))
        add_rect(s, 6.5, y, 0.06, 1.0, t["bg_dark"])
        add_text(s, cat, 6.7, y+0.08, 2.0, 0.35, size=13, bold=True,
                 color=t["text_dark"], font="Calibri")
        add_text(s, p, 6.7, y+0.42, 1.2, 0.35, size=20, bold=True,
                 color=t["bg_dark"], font="Calibri")
        add_text(s, amt, 8.0, y+0.42, 1.6, 0.35, size=12,
                 color=t["text_muted"], font="Calibri", align=PP_ALIGN.RIGHT)

    add_notes(s, "Revenue composition:\n• Food sales: 312 500 kr (62%) — driven by higher covers\n• Beverage: 175 000 kr (35%) — champagne sales strong\n• Catering event: 18 400 kr (3.6%) — one corporate event\n\nRevenue vs budget: +3.2% over March budget of 490 000 kr")
    return s

def slide_cost_breakdown(prs, t, d):
    """Slide 4: Cost breakdown with pie chart"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Cost Breakdown", 0.5, 0.17, 7, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")

    # Pie / doughnut chart
    chart_data = ChartData()
    chart_data.categories = ["Personal", "Råvaror & Dryck", "Lokal", "Övrigt"]
    chart_data.add_series("Costs", [d["staff"], d["food"], d["rent"], d["other"]])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(0.3), Inches(1.0), Inches(5.2), Inches(4.3),
        chart_data
    ).chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Color the segments
    colors_hex = [t["chart_colors"][0] if isinstance(t["chart_colors"][0],str) else "CADCFC",
                  "4A7FC1","FFC040","8AA8D0"]
    try:
        for i, pt in enumerate(chart.plots[0].series[0].points):
            pt.format.fill.solid()
            c = colors_hex[i % len(colors_hex)]
            pt.format.fill.fore_color.rgb = RGBColor(
                int(c[0:2],16), int(c[2:4],16), int(c[4:6],16))
    except Exception:
        pass

    # Right panel — cost rows with % bars
    costs = [
        ("Personal",        d["staff"],  d["staff_pct"],  40.0, "E24B4A"),
        ("Råvaror & dryck", d["food"],   d["food_pct"],   31.0, "1E2761"),
        ("Lokal & fastighet",d["rent"],  d["rent_pct"],   13.0, "4A7FC1"),
        ("Övrigt",          d["other"],  d["other_pct"],   6.0, "8AA8D0"),
    ]
    add_text(s, "vs target", 5.9, 1.05, 3.8, 0.38, size=11,
             color=t["text_muted"], font="Calibri")

    for i, (name, amt, actual_pct, target_pct, color_hex) in enumerate(costs):
        y = 1.5 + i * 0.97
        c = RGBColor(int(color_hex[:2],16),int(color_hex[2:4],16),int(color_hex[4:],16))
        add_text(s, name, 5.7, y, 2.4, 0.35, size=12, bold=True,
                 color=t["text_dark"], font="Calibri")
        add_text(s, f"{pct(actual_pct)}  (target {pct(target_pct)})", 8.1, y, 1.8, 0.35,
                 size=11, color=t["text_muted"], font="Calibri", align=PP_ALIGN.RIGHT)
        # Bar background
        add_rect(s, 5.7, y+0.38, 3.9, 0.18, rgb(0xE2,0xE8,0xF0))
        # Actual bar
        bar_w = min(3.9, 3.9 * float(actual_pct) / 70)
        add_rect(s, 5.7, y+0.38, bar_w, 0.18, c)
        # Target marker
        tm_x = 5.7 + 3.9 * float(target_pct) / 70
        add_rect(s, tm_x, y+0.32, 0.04, 0.3, rgb(0x60,0x60,0x60))
        add_text(s, kr(amt), 5.7, y+0.6, 3.9, 0.28, size=10,
                 color=t["text_muted"], font="Calibri")

    add_notes(s, "Cost analysis:\n• Staff 61.5% — well above 40% target. Root cause: sick leave weeks 11-12 required 42 000 kr extra overtime\n• Food/Bev 29.5% — within target. Menigo +8% partially offset by menu repricing\n• Rent fixed at 52 000 kr/month. Total including utilities 64 500 kr (12.7%)\n• Action: Board to approve staffing schedule review and Menigo renegotiation")
    return s

def slide_pl_statement(prs, t, d):
    """Slide 5: Full P&L table"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Profit & Loss Statement", 0.5, 0.17, 7, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")

    pl_rows = [
        ("INTÄKTER", "",  "", True, t["bg_dark"]),
        ("Matintäkter", kr(312500), "61.8%", False, None),
        ("Dryckesintäkter", kr(175000), "34.6%", False, None),
        ("Cateringintäkter", kr(18400), "3.6%", False, None),
        ("Totala intäkter", kr(d["revenue"]), "100%", True, rgb(0xE6,0xF1,0xFB)),
        ("KOSTNADER", "", "", True, t["bg_dark"]),
        ("Personalkostnader", kr(d["staff"]),  pct(d["staff_pct"]), False, None),
        ("Råvaror & dryck",   kr(d["food"]),   pct(d["food_pct"]),  False, None),
        ("Lokal & fastighet", kr(d["rent"]),   pct(d["rent_pct"]),  False, None),
        ("Övriga kostnader",  kr(d["other"]),  pct(d["other_pct"]), False, None),
        ("Totala kostnader",  kr(d["total_costs"]),
                              pct(round(d["total_costs"]/d["revenue"]*100,1)), True, rgb(0xFC,0xEB,0xEB)),
        ("RÖRELSERESULTAT (EBIT)", kr(d["profit"]),
                                    pct(d["margin"]), True, rgb(0xEA,0xF3,0xDE)),
    ]

    row_h = 0.36
    col_x = [0.4, 5.2, 7.8]
    col_w = [4.6, 2.4, 1.8]

    for i, (label, amount, percent, is_header, bg) in enumerate(pl_rows):
        y = 1.0 + i * row_h
        if bg:
            add_rect(s, 0.4, y, 9.2, row_h - 0.03, bg)
        text_color = t["white"] if bg == t["bg_dark"] else t["text_dark"]
        text_color_r = t["white"] if bg == t["bg_dark"] else rgb(0x3A,0x3A,0x3A)
        add_text(s, label,  col_x[0], y+0.04, col_w[0], row_h-0.08,
                 size=11 if not is_header else 11,
                 bold=is_header, color=text_color, font="Calibri")
        if amount:
            add_text(s, amount, col_x[1], y+0.04, col_w[1], row_h-0.08,
                     size=11, bold=is_header, color=text_color_r, font="Calibri",
                     align=PP_ALIGN.RIGHT)
            add_text(s, percent, col_x[2], y+0.04, col_w[2], row_h-0.08,
                     size=11, bold=is_header, color=text_color_r, font="Calibri",
                     align=PP_ALIGN.RIGHT)

    add_notes(s, "P&L highlights:\n• Top-line growth of 9.6% is excellent\n• Cost ratio at 92.3% vs 91.6% in February — margin compressed\n• EBIT of 38 700 kr held flat vs February despite higher revenue — all upside eaten by costs\n• Key action: Return cost ratio to below 90% in Q2")
    return s

def slide_profit_trend(prs, t, d):
    """Slide 6: Profit trend line chart + YTD comparison"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Profit Trend & Q1 Performance", 0.5, 0.17, 8, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")

    # Line chart — monthly profit
    chart_data = ChartData()
    chart_data.categories = d["monthly_labels"]
    chart_data.add_series("Net Profit (kr)", d["monthly_profit"])
    chart_data.add_series("Revenue (kr/10)", [x//10 for x in d["monthly_revenue"]])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.3), Inches(1.1), Inches(5.8), Inches(3.8),
        chart_data
    ).chart

    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Right panel — Q1 summary
    q1_items = [
        ("Q1 Revenue",  kr(d["q1_revenue"]),   kr(d["q1_budget"]), "+0.5%"),
        ("Q1 Costs",    kr(d["q1_revenue"] - 199500), "—",         "—"),
        ("Q1 Profit",   kr(199500),             kr(94500),          "+111%"),
        ("Avg Margin",  "14.7%",                "7.0%",             "+7.7pp"),
    ]
    add_text(s, "Q1 Actuals vs Budget", 6.3, 1.1, 3.5, 0.38,
             size=13, bold=True, color=t["text_dark"], font="Calibri")
    headers = ["Metric", "Actual", "Budget", "Δ"]
    col_xs = [6.3, 7.6, 8.5, 9.3]
    for j, h in enumerate(headers):
        add_text(s, h, col_xs[j], 1.5, 0.9, 0.3, size=9, bold=True,
                 color=t["text_muted"], font="Calibri")
    for i, (metric, actual, budget, delta) in enumerate(q1_items):
        y = 1.85 + i * 0.87
        add_rect(s, 6.3, y, 3.4, 0.75, rgb(0xFF,0xFF,0xFF))
        add_rect(s, 6.3, y, 0.05, 0.75, t["bg_dark"])
        add_text(s, metric, 6.42, y+0.08, 1.1, 0.28, size=10, bold=True,
                 color=t["text_dark"], font="Calibri")
        add_text(s, actual, col_xs[1], y+0.08, 0.85, 0.28, size=10, bold=True,
                 color=t["text_dark"], font="Calibri")
        add_text(s, budget, col_xs[2], y+0.08, 0.75, 0.28, size=10,
                 color=t["text_muted"], font="Calibri")
        delta_color = rgb(0x52,0xC3,0x74) if "+" in str(delta) else rgb(0xE2,0x4B,0x4A)
        add_text(s, delta, col_xs[3], y+0.08, 0.65, 0.28, size=10, bold=True,
                 color=delta_color, font="Calibri")

    add_notes(s, "Trend analysis:\n• Jan was below budget (slow start to year)\n• Feb and March recovered strongly\n• Q1 total revenue 1 356 900 kr vs budget 1 350 000 kr — slightly ahead\n• Profitability hampered by cost escalation; Q2 focus must be cost discipline")
    return s

def slide_staff_analysis(prs, t, d):
    """Slide 7: Staff cost deep-dive"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Staff Cost Analysis", 0.5, 0.17, 7, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")
    add_text(s, "⚠  Requires action", 7.5, 0.22, 2.2, 0.42,
             size=12, bold=True, color=rgb(0xFF,0xC0,0x40), font="Calibri",
             align=PP_ALIGN.RIGHT)

    # Left: staff breakdown bar chart
    chart_data = ChartData()
    chart_data.categories = ["Kök", "Servering", "Disk", "Övertid"]
    chart_data.add_series("Lön + arb.avg (kr)", [130000, 58000, 13000, 23978])

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.3), Inches(1.1), Inches(5.5), Inches(3.8),
        chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    for pt in plot.series[0].points:
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = t["bg_dark"]

    # Right: 3 action cards
    actions = [
        ("Immediate", "Review week 11–12 overtime billing.\nVerify sick leave documentation."),
        ("This month", "Rebuild staffing schedule to reduce overtime dependency.\nTarget: 40% of revenue."),
        ("Q2", "Recruit permanent kock to reduce reliance on temp staff.\nProject savings: ~28 000 kr/month."),
    ]
    colors_a = [rgb(0xE2,0x4B,0x4A), rgb(0xFF,0xC0,0x40), t["bg_dark"]]
    for i, (timing, action) in enumerate(actions):
        y = 1.1 + i * 1.5
        add_rect(s, 5.95, y, 3.8, 1.35, rgb(0xFF,0xFF,0xFF))
        add_rect(s, 5.95, y, 0.08, 1.35, colors_a[i])
        add_text(s, timing.upper(), 6.15, y+0.1, 3.5, 0.35,
                 size=10, bold=True, color=colors_a[i], font="Calibri")
        add_text(s, action, 6.15, y+0.45, 3.4, 0.8,
                 size=11, color=t["text_dark"], font="Calibri")

    add_notes(s, "Staff cost drivers:\n• Kitchen (köksmästare + 3 kockar): 173 143 kr\n• Service (4 servitriser + 2 timanställda): 87 181 kr\n• Dishwashers: 24 881 kr\n• Overtime (sick leave cover): 42 000 kr EXCESS — this is the primary problem\n\nBoard decision needed: Approve recruitment of permanent kock (target start Q3).")
    return s

def slide_recommendations(prs, t, d):
    """Slide 8: Recommendations with priority flags"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Recommendations & Actions", 0.5, 0.17, 7, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")

    recs = [
        ("HIGH",   "Renegotiate Menigo contract",
         "March +8% price increase adds ~15 000 kr/month. Request Q2 price freeze or trial Martin & Servera for 30% of volume. Potential saving: 5 000–8 000 kr/month."),
        ("HIGH",   "Staffing schedule review",
         "Overtime at 42 000 kr caused by sick-leave cover. Implement cross-training and on-call pool. Target: staff costs ≤40% of revenue."),
        ("MEDIUM", "Q2 budget revision",
         "Current Q2 budget assumes old cost levels. Revise food cost budget +8% and staff budget +5% based on March actuals."),
        ("MEDIUM", "Cash flow buffer",
         "Bank statement shows negative intra-month balance around payroll (25th). Maintain ≥150 000 kr minimum working capital."),
        ("LOW",    "POS data quality",
         "Ancon POS Fortnox sync shows 3 uncategorised transactions. Clean up before Q2 to improve reporting accuracy."),
    ]
    colors_p = {"HIGH": rgb(0xE2,0x4B,0x4A), "MEDIUM": rgb(0xFF,0xC0,0x40), "LOW": rgb(0x52,0xC3,0x74)}

    for i, (priority, title, detail) in enumerate(recs):
        y = 1.0 + i * 0.92
        c = colors_p[priority]
        # Priority chip
        add_rect(s, 0.4, y+0.12, 0.85, 0.32, c)
        add_text(s, priority, 0.42, y+0.12, 0.81, 0.32,
                 size=9, bold=True, color=rgb(0xFF,0xFF,0xFF),
                 font="Calibri", align=PP_ALIGN.CENTER)
        # Title
        add_text(s, title, 1.38, y+0.06, 3.2, 0.38,
                 size=13, bold=True, color=t["text_dark"], font="Calibri")
        # Detail
        add_text(s, detail, 4.7, y+0.04, 5.0, 0.78,
                 size=10, color=rgb(0x55,0x55,0x55), font="Calibri")
        # Divider
        if i < len(recs)-1:
            add_rect(s, 0.4, y+0.88, 9.2, 0.02, rgb(0xE2,0xE8,0xF0))

    add_notes(s, "Board decisions required:\n1. Approve Menigo renegotiation mandate (owner to lead)\n2. Approve permanent kock recruitment budget\n3. Note Q2 budget revision (CFO/owner to update)\n4. Acknowledge cash flow monitoring requirement\n\nAll HIGH items require owner action within 2 weeks.")
    return s

def slide_next_steps(prs, t, d):
    """Slide 9: Next steps / action items"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, rgb(0xF8,0xF8,0xFA))

    add_rect(s, 0, 0, 10, 0.9, t["bg_dark"])
    add_text(s, "Next Steps", 0.5, 0.17, 7, 0.56,
             size=22, bold=True, color=t["white"], font="Calibri")

    timeline = [
        ("Week 1",  "Contact Menigo for pricing meeting",      "Owner"),
        ("Week 1",  "Pull overtime data — weeks 11–12",        "Köksmästare"),
        ("Week 2",  "Rebuild staffing schedule model",         "Owner + Köksmästare"),
        ("Week 2",  "Revise Q2 budget figures",                "Owner"),
        ("Month end","Verify Ancon–Fortnox sync accuracy",     "Admin"),
        ("Q2",      "Launch permanent kock recruitment",       "Owner"),
    ]

    col_xs = [0.4, 1.8, 5.5, 8.8]
    headers = ["When", "Action", "Owner", ""]
    for j, h in enumerate(headers[:3]):
        add_text(s, h.upper(), col_xs[j], 1.0, 1.2, 0.35,
                 size=9, bold=True, color=t["text_muted"], font="Calibri")

    for i, (when, action, owner) in enumerate(timeline):
        y = 1.4 + i * 0.67
        bg = rgb(0xFF,0xFF,0xFF) if i % 2 == 0 else rgb(0xF4,0xF6,0xFA)
        add_rect(s, 0.4, y, 9.2, 0.6, bg)
        # When chip
        add_rect(s, 0.42, y+0.14, 1.15, 0.3, t["bg_dark"])
        add_text(s, when, 0.44, y+0.14, 1.11, 0.3,
                 size=9, bold=True, color=t["white"],
                 font="Calibri", align=PP_ALIGN.CENTER)
        add_text(s, action, 1.75, y+0.14, 6.8, 0.3,
                 size=12, color=rgb(0x2A,0x2A,0x2A), font="Calibri")
        add_text(s, owner, 8.7, y+0.14, 1.1, 0.3,
                 size=11, color=t["text_muted"], font="Calibri", align=PP_ALIGN.RIGHT)

    add_notes(s, "Close meeting by confirming:\n• Owner to schedule Menigo meeting this week\n• Next monthly review: first week of May\n• Interim check-in on staffing: April 17")
    return s

def slide_closing(prs, t, d):
    """Slide 10: Closing / thank you slide"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, t["bg_mid"])

    add_rect(s, 0, 0, 10, 0.08, t["accent"])
    add_rect(s, 0, 5.545, 10, 0.08, t["accent"])

    add_text(s, "Thank you", 1.0, 1.8, 8, 1.2,
             size=48, bold=True, color=t["white"], font="Calibri",
             align=PP_ALIGN.CENTER)
    add_text(s, f"Questions & Discussion", 1.0, 3.1, 8, 0.6,
             size=20, color=t["accent"], font="Calibri",
             align=PP_ALIGN.CENTER, italic=True)
    add_text(s, f"{d['company']}  ·  {d['period']}  ·  Konfidentiellt",
             1.0, 4.8, 8, 0.4,
             size=11, color=t["text_muted"], font="Calibri",
             align=PP_ALIGN.CENTER)

    add_notes(s, "Open for questions.\n\nKey reminders for follow-up:\n• Send action list to all attendees within 24 hours\n• Next meeting: first week of May")
    return s

# ════════════════════════════════════════════════════════════════
# DECK ASSEMBLERS — combine slides into full presentations
# ════════════════════════════════════════════════════════════════

def build_deck(deck_type="board", theme_key="midnight", data_overrides=None,
               output_path="/home/claude/studio/presentation.pptx"):
    d = get_data(data_overrides)
    t = THEMES.get(theme_key, THEMES["midnight"])

    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    deck_configs = {
        "board": {
            "title":    f"Board Meeting — {d['period']}",
            "subtitle": f"Monthly Performance Review",
            "type":     "Board Meeting",
            "slides":   [slide_title, slide_exec_summary, slide_revenue,
                         slide_cost_breakdown, slide_pl_statement, slide_profit_trend,
                         slide_staff_analysis, slide_recommendations,
                         slide_next_steps, slide_closing],
        },
        "investor": {
            "title":    f"Investor Update — {d['quarter']}",
            "subtitle": "Business Performance & Growth Strategy",
            "type":     "Investor Pitch",
            "slides":   [slide_title, slide_exec_summary, slide_revenue,
                         slide_cost_breakdown, slide_pl_statement, slide_profit_trend,
                         slide_recommendations, slide_closing],
        },
        "monthly": {
            "title":    f"Monthly Review — {d['period']}",
            "subtitle": "Operations & Financial Summary",
            "type":     "Monthly Review",
            "slides":   [slide_title, slide_exec_summary, slide_revenue,
                         slide_cost_breakdown, slide_profit_trend,
                         slide_recommendations, slide_next_steps, slide_closing],
        },
        "budget": {
            "title":    f"Budget Proposal — Q2 2026",
            "subtitle": "Cost Revisions & Forward Projections",
            "type":     "Budget Proposal",
            "slides":   [slide_title, slide_exec_summary, slide_cost_breakdown,
                         slide_staff_analysis, slide_profit_trend,
                         slide_recommendations, slide_next_steps, slide_closing],
        },
    }

    cfg = deck_configs.get(deck_type, deck_configs["board"])

    # Build all slides
    for i, builder in enumerate(cfg["slides"]):
        if builder == slide_title:
            builder(prs, t, d, cfg["title"], cfg["subtitle"], cfg["type"])
        else:
            builder(prs, t, d)

    prs.save(output_path)
    return output_path, len(cfg["slides"])

# ════════════════════════════════════════════════════════════════
# CLI entry point
# ════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--deck",   default="board",    help="board|investor|monthly|budget")
    parser.add_argument("--theme",  default="midnight", help="midnight|forest|terracotta")
    parser.add_argument("--output", default="/home/claude/studio/presentation.pptx")
    parser.add_argument("--data",   default=None,       help="JSON string of data overrides")
    args = parser.parse_args()

    overrides = json.loads(args.data) if args.data else None
    path, n = build_deck(args.deck, args.theme, overrides, args.output)
    print(json.dumps({"status": "ok", "path": path, "slides": n,
                      "deck": args.deck, "theme": args.theme}))
